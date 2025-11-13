#!/usr/bin/env python3

import sys
from pathlib import Path
from pptx import Presentation

if len(sys.argv) < 2:
    print("Usage: python pptx_to_md.py <input.pptx>")
    sys.exit(1)

input_file = sys.argv[1]
output_file = Path(input_file).with_suffix('.md')
presentation = Presentation(input_file)

with open(output_file, 'w', encoding='utf-8') as f:
    for i, slide in enumerate(presentation.slides, 1):
        f.write(f"## Slide {i}\n\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                f.write(shape.text + "\n\n")

print(f"Converted {input_file} to {output_file}")
